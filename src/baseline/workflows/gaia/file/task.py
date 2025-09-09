def task1_file_process(args, dag_id, question, supplementary_files, vllm_manager= None):
    """
    Processes the file content based on type and stores the processed content in context.
    """
    import time
    import os
    from docx import Document
    from pptx import Presentation
    import xml.etree.ElementTree as ET
    from xml.dom import minidom
    import pandas as pd
    from io import BytesIO    
    import pdfplumber  # 添加 PDF 处理库    
    try:
        start_time= time.time()
        file_name = ""
        content = None
        
        if supplementary_files:
            file_name = next(iter(supplementary_files.keys()), "")
            content = next(iter(supplementary_files.values()), None)
            print(f"  -> Found single supplementary file to process: '{file_name}'")
        else:
            print("  -> No supplementary files found in context.")
        processed_content = ""
        
        print("✅ Task 1: Starting file processing...")
        if content is None:
            print(f"  -> No content for file '{file_name}', skipping processing.")
        else:
            print(f"  -> Processing file: {file_name}")
            file_name_lower = file_name.lower()
            
            # Excel 文件处理
            if file_name_lower.endswith(('.xlsx', '.xls')):
                print("处理 Excel 文件...")
                excel_file = BytesIO(content)
                
                # 读取所有sheet
                excel_data = []
                with pd.ExcelFile(excel_file) as xls:
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(xls, sheet_name)
                        excel_data.append(f"=== Sheet: {sheet_name} ===\n")
                        excel_data.append(df.to_string(index=True))
                
                processed_content = "\n\n".join(excel_data)
                print(f"提取到 {len(excel_data)//2} 个工作表的内容")

            # CSV 文件处理
            elif file_name_lower.endswith('.csv'):
                print("处理 CSV 文件...")
                csv_file = BytesIO(content)
                df = pd.read_csv(csv_file)
                processed_content = df.to_string(index=True)
                print(f"提取到 CSV 内容，共 {len(df)} 行")

            # XML 文件处理
            elif file_name_lower.endswith('.xml'):
                print("[File Process] 处理 XML 文件...")
                try:
                    xml_content = content.decode('utf-8')
                    print("[File Process] XML 内容解码成功")
                    
                    # 尝试解析和格式化 XML
                    try:
                        parsed_xml = minidom.parseString(xml_content)
                        processed_content = parsed_xml.toprettyxml()
                        print("[File Process] XML 格式化成功")
                    except Exception as xml_error:
                        print(f"[File Process] XML 格式化失败: {str(xml_error)}")
                        processed_content = xml_content
                    
                    # 提取文本内容（去除 XML 标签）
                    try:
                        root = ET.fromstring(xml_content)
                        def extract_text(element):
                            text = element.text or ''
                            for child in element:
                                text += extract_text(child)
                                if child.tail:
                                    text += child.tail
                            return text
                        
                        text_content = extract_text(root)
                        print("[File Process] 文本内容提取成功")
                        
                        processed_content = (
                            "=== XML Structure ===\n"
                            f"{processed_content}\n\n"
                            "=== Extracted Text Content ===\n"
                            f"{text_content}"
                        )
                    except Exception as text_error:
                        print(f"[File Process] 文本提取失败: {str(text_error)}")
                    
                    print(f"[File Process] XML 处理完成，内容长度: {len(processed_content)}")
                    
                except Exception as e:
                    print(f"[File Process] XML 处理出错: {str(e)}")
                    processed_content = content.decode('utf-8')

            # PPT 文件处理
            elif file_name_lower.endswith('.pptx'):
                print("处理 PPT 文件...")
                prs = Presentation(BytesIO(content))
                all_slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text_parts.append(shape.text.strip())
                    all_slides_text.append(f"--- Slide {i+1} ---\n" + '\n'.join(filter(None, slide_text_parts)))
                processed_content = "\n\n".join(all_slides_text)
                print(f"提取到 {len(all_slides_text)} 页幻灯片内容")

            # Word 文件处理
            elif file_name_lower.endswith('.docx'):
                print("处理 Word 文件...")
                doc = Document(BytesIO(content))
                paragraphs_text = []
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs_text.append(text)
                processed_content = "\n".join(paragraphs_text)
                print(f"提取到 {len(paragraphs_text)} 段文本")

            # PDF 文件处理
            elif file_name_lower.endswith('.pdf'):
                print("处理 PDF 文件...")
                temp_pdf_path = f"/tmp/temp_{file_name_lower.split('.')[0]}.pdf"
                with open(temp_pdf_path, "wb") as f:
                    f.write(content)
                
                pdf_pages = []
                with pdfplumber.open(temp_pdf_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            pdf_pages.append(f"--- Page {i+1} ---\n{page_text}")
                        
                        tables = page.extract_tables()
                        if tables:
                            table_texts = []
                            for j, table in enumerate(tables):
                                table_text = f"Table {j+1} on Page {i+1}:\n"
                                for row in table:
                                    cleaned_row = [str(cell) if cell is not None else "" for cell in row]
                                    table_text += " | ".join(cleaned_row) + "\n"
                                table_texts.append(table_text)
                            
                            if table_texts:
                                pdf_pages.append("\n".join(table_texts))
                
                os.remove(temp_pdf_path)
                processed_content = "\n\n".join(pdf_pages)
                print(f"提取到 {len(pdf_pages)} 页/节 PDF 内容")

            # 文本文件处理
            elif file_name_lower.endswith('.txt'):
                print("处理 TXT 文件...")
                processed_content = content.decode('utf-8')

            else:
                print(f"警告: 未知或不支持的文件类型: {file_name_lower}，尝试按文本处理")
                try:
                    processed_content = content.decode('utf-8')
                except UnicodeDecodeError:
                    print("文件无法解码为 UTF-8，将跳过内容注入")
                    processed_content = "[Content could not be decoded]"
            print("✅ Task 1: File processing finished.")

        print(f"✅ 文件处理结束...")
        return {
            "dag_id": dag_id, 
            "question": question,
            "processed_content": processed_content,
            "start_time": start_time,
            "end_time": time.time()
        }
    except Exception as e:
        print(f"task1_file_process 发生错误: {str(e)}")
        raise e

def task2_llm_process_qwen(args, dag_id, question, processed_content, vllm_manager= None, backend= "huggingface"):
    """
    Processes the file content using LLM based on the question.
    """
    import gc
    import torch
    from typing import List, Dict
    def query_vllm_model(api_url: str, model_alias: str, messages: List, temperature: float= 0.6, max_token: int= 1024, top_p: float= 0.9, repetition_penalty: float= 1.1):
        """
        通过HTTP请求查询本地vLLM服务。
        
        :param api_url: vLLM服务的根URL (例如 http://127.0.0.1:8000)
        :param model_alias: 在vLLM中服务的模型别名 (例如 qwen3-32b)
        :param messages: OpenAI格式的消息列表
        :param temperature: 温度参数
        :param max_token: 最大生成token数
        :param top_p: Top-p采样参数
        :param repetition_penalty: 重复惩罚参数
        :return: 一个包含性能特征和模型响应文本的元组
        """
        import requests
        from rich.console import Console
        console = Console()
        # vLLM的聊天接口路径
        chat_url = f"{api_url.strip('/')}/v1/chat/completions"
        headers = {
            "Content-Type": "application/json"
        }
        # 构建与OpenAI兼容的请求体
        payload = {
            "model": model_alias,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_token,
            "top_p": top_p,
            "repetition_penalty": repetition_penalty,
        }
        # 为了计算特征，先将消息拼接起来
        conversation = ""
        for message in messages:
            conversation += f"{message['role']}: {message['content']}"
        try:
            response = requests.post(chat_url, json=payload, headers=headers, timeout=3600)
            response.raise_for_status()  # 如果请求失败则抛出异常
            response_data = response.json()
            content = response_data['choices'][0]['message']['content'].lstrip()
            return content

        except requests.exceptions.RequestException as e:
            error_msg = f"vLLM request failed: {str(e)}"
            return f"[bold red]{error_msg}"
    def query_llm(model, model_folder, messages, temperature= 0.6, max_token= 1024, top_p= 0.9, repetition_penalty= 1.1):
        import os
        from transformers import AutoTokenizer, AutoModelForCausalLM
        
        # 加载预训练模型和分词器
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        tokenizer = AutoTokenizer.from_pretrained(tokenizer_path, device_map= "cuda")
        local_model = AutoModelForCausalLM.from_pretrained(
            os.path.join(model_folder, model),
            torch_dtype="float16",
            device_map="cuda", offload_state_dict= False,
            low_cpu_mem_usage=True
            )

        # 将messages转换成字符串格式
        conversation = ""
        for message in messages:
            conversation += f"{message['role']}: {message['content']}"

        # 编码输入
        input_ids = tokenizer.encode(conversation + tokenizer.eos_token, return_tensors='pt').to("cuda")

        # 生成回复
        output = local_model.generate(
            input_ids, 
            pad_token_id=tokenizer.eos_token_id,
            max_new_tokens= max_token,
            temperature= temperature, 
            top_p= top_p,
            repetition_penalty= repetition_penalty
        )
        response = tokenizer.decode(output[:, input_ids.shape[-1]:][0], skip_special_tokens=True)

        del local_model
        del tokenizer
        del input_ids
        del output
        gc.collect()
        torch.cuda.empty_cache()
        if torch.cuda.is_available():
            torch.cuda.synchronize()
            torch.cuda.reset_max_memory_allocated()
            torch.cuda.reset_peak_memory_stats()
        return response
    
    def query_llm_online(api_url, api_key, payload)-> str:
        """Call SiliconFlow API to get LLM response"""
        import requests
        from rich.console import Console

        console= Console()
        headers= {
            "Authorization": api_key,
            "Content-Type": "application/json"
        }
        try:
            response= requests.post(
                api_url,
                json= payload,
                headers= headers,
                timeout= 3600
            )
            response.raise_for_status()
            return response.json()['choices'][0]['message']['content'].lstrip()
        except Exception as e:
            console.print(f"[bold red]API request failed: {str(e)}")
            return f"[bold red]API request failed: {str(e)}"
    import time
    try:
        start_time= time.time()
        print(f"✅ qwen处理开始....")
        if not question:
             raise ValueError(f"任务 {dag_id} 缺少 Question")

        # 构建提示
        # 生成推理模型的特征
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        payload= {
            "model": "Qwen/Qwen3-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": args.temperature,
            "max_tokens": args.max_token,
            "enable_thinking": False,
            "response_format": {"type": "text"}               
        }
        answer= None
        if args.use_online_model:
            answer= query_llm_online(
                api_url= args.api_url,
                api_key= args.api_key, 
                payload= payload)
        elif backend== "vllm":
            vllm_api_url= vllm_manager.get_next_endpoint("qwen3-32b")
            answer= query_vllm_model(api_url= vllm_api_url, model_alias= "qwen3-32b", messages= [{"role": "user", "content": prompt}], temperature= args.temperature, max_token= args.max_token, top_p= args.top_p, repetition_penalty= args.repetition_penalty)
        else:
            answer= query_llm(model= "Qwen/Qwen3-32B", model_folder= args.model_folder, messages= [{"role": "user", "content": prompt}], temperature= args.temperature, max_token= args.max_token, top_p= args.top_p, repetition_penalty= args.repetition_penalty)
        return {
            "dag_id": dag_id,
            "question": question,
            "answer_qwen": answer,
            "start_time": start_time,
            "end_time": time.time()            
        }
    except Exception as e:
        print(f"task2_llm_process_qwen 发生错误: {str(e)}")
        raise e

def task3_llm_process_deepseek(args, dag_id, question, processed_content, vllm_manager= None, backend= "huggingface"):
    """
    Processes the file content using LLM based on the question.
    """
    import time
    from typing import List, Dict
    def query_vllm_model(api_url: str, model_alias: str, messages: List, temperature: float= 0.6, max_token: int= 1024, top_p: float= 0.9, repetition_penalty: float= 1.1):
        """
        通过HTTP请求查询本地vLLM服务。
        
        :param api_url: vLLM服务的根URL (例如 http://127.0.0.1:8000)
        :param model_alias: 在vLLM中服务的模型别名 (例如 qwen3-32b)
        :param messages: OpenAI格式的消息列表
        :param temperature: 温度参数
        :param max_token: 最大生成token数
        :param top_p: Top-p采样参数
        :param repetition_penalty: 重复惩罚参数
        :return: 一个包含性能特征和模型响应文本的元组
        """
        import requests
        from rich.console import Console
        console = Console()
        # vLLM的聊天接口路径
        chat_url = f"{api_url.strip('/')}/v1/chat/completions"
        headers = {
            "Content-Type": "application/json"
        }
        # 构建与OpenAI兼容的请求体
        payload = {
            "model": model_alias,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_token,
            "top_p": top_p,
            "repetition_penalty": repetition_penalty,
        }
        # 为了计算特征，先将消息拼接起来
        conversation = ""
        for message in messages:
            conversation += f"{message['role']}: {message['content']}"
        try:
            response = requests.post(chat_url, json=payload, headers=headers, timeout=3600)
            response.raise_for_status()  # 如果请求失败则抛出异常
            response_data = response.json()
            content = response_data['choices'][0]['message']['content'].lstrip()
            return content

        except requests.exceptions.RequestException as e:
            error_msg = f"vLLM request failed: {str(e)}"
            return f"[bold red]{error_msg}"
    def query_llm(model, model_folder, messages, temperature= 0.6, max_token= 1024, top_p= 0.9, repetition_penalty= 1.1):
        import os
        import gc
        import torch
        from transformers import AutoTokenizer, AutoModelForCausalLM
        
        # 加载预训练模型和分词器
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        tokenizer = AutoTokenizer.from_pretrained(tokenizer_path, device_map= "cuda")
        local_model = AutoModelForCausalLM.from_pretrained(
            os.path.join(model_folder, model),
            torch_dtype="float16",
            device_map="cuda", offload_state_dict= False,
            low_cpu_mem_usage=True
            )

        # 将messages转换成字符串格式
        conversation = ""
        for message in messages:
            conversation += f"{message['role']}: {message['content']}"

        # 编码输入
        input_ids = tokenizer.encode(conversation + tokenizer.eos_token, return_tensors='pt').to("cuda")

        # 生成回复
        output = local_model.generate(
            input_ids, 
            pad_token_id=tokenizer.eos_token_id,
            max_new_tokens= max_token,
            temperature= temperature, 
            top_p= top_p,
            repetition_penalty= repetition_penalty
        )
        response = tokenizer.decode(output[:, input_ids.shape[-1]:][0], skip_special_tokens=True)

        del local_model
        del tokenizer
        del input_ids
        del output
        torch.cuda.empty_cache()
        gc.collect()
        # 强制释放所有CUDA资源
        if torch.cuda.is_available():
            torch.cuda.synchronize()
            torch.cuda.reset_max_memory_allocated()
            torch.cuda.reset_peak_memory_stats()
        return response
    
    def query_llm_online(api_url, api_key, payload)-> str:
        """Call SiliconFlow API to get LLM response"""
        import requests
        from rich.console import Console

        console= Console()
        headers= {
            "Authorization": api_key,
            "Content-Type": "application/json"
        }
        try:
            response= requests.post(
                api_url,
                json= payload,
                headers= headers,
                timeout= 3600
            )
            response.raise_for_status()
            return response.json()['choices'][0]['message']['content'].lstrip()
        except Exception as e:
            console.print(f"[bold red]API request failed: {str(e)}")
            return f"[bold red]API request failed: {str(e)}"    
    try:
        start_time= time.time()
        print(f"✅ deepseek处理开始....")
        if not question:
             raise ValueError(f"任务 {dag_id} 缺少 Question")

        # 构建提示
        # 生成推理模型的特征
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your concise thinking thoughts and don't think too complicated, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        payload= {
            "model": "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": args.temperature,
            "max_tokens": args.max_token,
            "response_format": {"type": "text"}               
        }
        answer= None
        if args.use_online_model:
            answer= query_llm_online(
                api_url= args.api_url,
                api_key= args.api_key, 
                payload= payload)
        elif backend== "vllm":
            vllm_api_url= vllm_manager.get_next_endpoint("deepseek-r1-32b")
            answer= query_vllm_model(api_url= vllm_api_url, model_alias= "deepseek-r1-32b", messages= [{"role": "user", "content": prompt}], temperature= args.temperature, max_token= args.max_token, top_p= args.top_p, repetition_penalty= args.repetition_penalty)
        else:
            answer= query_llm(model= "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B", model_folder= args.model_folder, messages= [{"role": "user", "content": prompt}], temperature= args.temperature, max_token= args.max_token, top_p= args.top_p, repetition_penalty= args.repetition_penalty)
        return {
            "dag_id": dag_id,
            "question": question,
            "answer_deepseek": answer,
            "start_time": start_time,
            "end_time": time.time()
        }
    except Exception as e:
        print(f"task3_llm_process_deepseek 发生错误: {str(e)}")
        raise e

def task4_fuse_llm_answer(args, dag_id, question, answer_qwen, answer_deepseek, vllm_manager= None, backend= "huggingface"):
    """
    Fusion Task: Takes answers from multiple experts, analyzes them, and generates a final, synthesized answer.
    """
    import time
    import torch, gc
    from typing import List, Dict
    def query_vllm_model(api_url: str, model_alias: str, messages: List, temperature: float= 0.6, max_token: int= 1024, top_p: float= 0.9, repetition_penalty: float= 1.1):
        """
        通过HTTP请求查询本地vLLM服务。
        
        :param api_url: vLLM服务的根URL (例如 http://127.0.0.1:8000)
        :param model_alias: 在vLLM中服务的模型别名 (例如 qwen3-32b)
        :param messages: OpenAI格式的消息列表
        :param temperature: 温度参数
        :param max_token: 最大生成token数
        :param top_p: Top-p采样参数
        :param repetition_penalty: 重复惩罚参数
        :return: 一个包含性能特征和模型响应文本的元组
        """
        import requests
        from rich.console import Console
        console = Console()
        # vLLM的聊天接口路径
        chat_url = f"{api_url.strip('/')}/v1/chat/completions"
        headers = {
            "Content-Type": "application/json"
        }
        # 构建与OpenAI兼容的请求体
        payload = {
            "model": model_alias,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_token,
            "top_p": top_p,
            "repetition_penalty": repetition_penalty,
        }
        # 为了计算特征，先将消息拼接起来
        conversation = ""
        for message in messages:
            conversation += f"{message['role']}: {message['content']}"
        try:
            response = requests.post(chat_url, json=payload, headers=headers, timeout=3600)
            response.raise_for_status()  # 如果请求失败则抛出异常
            response_data = response.json()
            content = response_data['choices'][0]['message']['content'].lstrip()
            return content

        except requests.exceptions.RequestException as e:
            error_msg = f"vLLM request failed: {str(e)}"
            return f"[bold red]{error_msg}"
    def query_llm(model, model_folder, messages, temperature= 0.6, max_token= 1024, top_p= 0.9, repetition_penalty= 1.1):
        import os
        from transformers import AutoTokenizer, AutoModelForCausalLM
        
        # 加载预训练模型和分词器
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        tokenizer = AutoTokenizer.from_pretrained(tokenizer_path, device_map= "cuda")
        local_model = AutoModelForCausalLM.from_pretrained(
            os.path.join(model_folder, model),
            torch_dtype="float16",
            device_map="cuda", offload_state_dict= False,
            low_cpu_mem_usage=True
            )

        # 将messages转换成字符串格式
        conversation = ""
        for message in messages:
            conversation += f"{message['role']}: {message['content']}"

        # 编码输入
        input_ids = tokenizer.encode(conversation + tokenizer.eos_token, return_tensors='pt').to("cuda")

        # 生成回复
        output = local_model.generate(
            input_ids, 
            pad_token_id=tokenizer.eos_token_id,
            max_new_tokens= max_token,
            temperature= temperature, 
            top_p= top_p,
            repetition_penalty= repetition_penalty
        )
        response = tokenizer.decode(output[:, input_ids.shape[-1]:][0], skip_special_tokens=True)

        del local_model
        del tokenizer
        del input_ids
        del output
        torch.cuda.empty_cache()
        gc.collect()
        if torch.cuda.is_available():
            torch.cuda.synchronize()
            torch.cuda.reset_max_memory_allocated()
            torch.cuda.reset_peak_memory_stats()
        return response
      
    def query_llm_online(api_url, api_key, payload)-> str:
        """Call SiliconFlow API to get LLM response"""
        import requests
        from rich.console import Console

        console= Console()
        headers= {
            "Authorization": api_key,
            "Content-Type": "application/json"
        }
        try:
            response= requests.post(
                api_url,
                json= payload,
                headers= headers,
                timeout= 3600
            )
            response.raise_for_status()
            return response.json()['choices'][0]['message']['content'].lstrip()
        except Exception as e:
            console.print(f"[bold red]API request failed: {str(e)}")
            return f"[bold red]API request failed: {str(e)}"
    try:
        start_time= time.time()
        print("✅ Task 4 (Fusion): Starting answer synthesis...")
        # 2. 构建一个用于融合的 prompt
        prompt = (
            "You are a senior editor and a world-class reasoning expert. Your job is to synthesize the answers from two different AI assistants to produce one final, superior answer for the given question.\n\n"
            f"--- Original Question ---\n{question}\n\n"
            f"--- Answer from Assistant 1 (Qwen3) ---\n{answer_qwen}\n\n"
            f"--- Answer from Assistant 2 (DeepSeek) ---\n{answer_deepseek}\n\n"
            "--- Your Task ---\n"
            "Analyze both answers. Identify the strengths and weaknesses of each. Then, combine their best elements, correct any errors, and provide a single, comprehensive, and accurate final answer. Adhere to the final answer format requested in the original prompt.\n\n"
            "Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]."
        )
        
        # 3. 使用一个强大的模型（如Qwen3）来执行融合
        payload= {
            "model": "Qwen/Qwen3-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": args.temperature,
            "max_tokens": args.max_token,
            "response_format": {"type": "text"},            
            "enable_thinking": False
        }
        final_answer= None
        if args.use_online_model:
            final_answer= query_llm_online(
                api_url= args.api_url,
                api_key= args.api_key, 
                payload= payload)
        elif backend== "vllm":
            vllm_api_url= vllm_manager.get_next_endpoint("qwen3-32b")
            final_answer= query_vllm_model(api_url= vllm_api_url, model_alias= "qwen3-32b", messages= [{"role": "user", "content": prompt}], temperature= args.temperature, max_token= args.max_token, top_p= args.top_p, repetition_penalty= args.repetition_penalty)
        else:
            final_answer= query_llm(model= "Qwen/Qwen3-32B", model_folder= args.model_folder, messages= [{"role": "user", "content": prompt}], temperature= args.temperature, max_token= args.max_token, top_p= args.top_p, repetition_penalty= args.repetition_penalty)
        # 4. 准备并返回最终的输出
        print(f"  -> Fusion complete. Final answer generated for DAG {dag_id}.")
        return {
            "dag_id": dag_id,
            "final_answer": final_answer,
            "start_time": start_time,
            "end_time": time.time()
        }
    except Exception as e:
        print(f"❌ task4_fuse_final_answer failed: {str(e)}")
        raise e