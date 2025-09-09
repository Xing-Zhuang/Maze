import os
import ray
import time
import gc
import torch
import json
import pdfplumber  # 添加 PDF 处理库
import pandas as pd
from io import BytesIO
from docx import Document
from pptx import Presentation
import xml.etree.ElementTree as ET
from xml.dom import minidom
from typing import List, Dict
from agentos.scheduler import cpu,gpu,io

import re
def estimate_tokens(text):
    """
    一个更精确的、用于估算中英混合文本token数的函数。
    它分别计算CJK字符和非CJK单词，并避免了重复计算。
    """
    # 1. 精确计算CJK字符数
    cjk_chars = sum(1 for char in text if '\u4E00' <= char <= '\u9FFF')
    # 2. 从文本中移除所有CJK字符和换行符，以便统计其他语言的单词
    # 我们用正则表达式将CJK字符替换为空格，以确保单词能被正确分割
    non_cjk_text = re.sub(r'[\u4E00-\u9FFF]', ' ', text)
    non_cjk_text = non_cjk_text.replace("\n", " ")
    # 3. 计算非CJK单词数，乘以经验系数
    # 使用split()可以有效地按空格分割出单词
    non_cjk_words_count = len(non_cjk_text.split())
    # 4. 将两部分加总
    # 这里的1.3是英文单词到token的经验系数，可以微调
    estimated_tokens = cjk_chars + int(non_cjk_words_count * 1.3)
    return estimated_tokens

def query_vllm_model(api_url: str, model_alias: str, messages: List, temperature: float= 0.6, max_token: int= 1024, top_p: float= 0.9, repetition_penalty: float= 1.1) -> tuple[dict, str]:
    """
    通过HTTP请求查询本地vLLM服务。
    
    :param api_url: vLLM服务的根URL (例如 http://127.0.0.1:8000)
    :param model_alias: 在vLLM中服务的模型别名 (例如 qwen3-32b)
    :param messages: OpenAI格式的消息列表
    :param temperature: 温度参数
    :param max_token: 最大生成token数
    :param top_p: Top-p采样参数
    :param repetition_penalty: 重复惩罚参数
    :return: 一个包含性能特征和模型响应文本的元组
    """
    import requests
    from rich.console import Console
    console = Console()
    # vLLM的聊天接口路径
    chat_url = f"{api_url.strip('/')}/v1/chat/completions"
    headers = {
        "Content-Type": "application/json"
    }
    # 构建与OpenAI兼容的请求体
    payload = {
        "model": model_alias,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_token,
        "top_p": top_p,
        "repetition_penalty": repetition_penalty,
    }
    # 为了计算特征，先将消息拼接起来
    conversation = ""
    for message in messages:
        conversation += f"{message['role']}: {message['content']}"
    try:
        response = requests.post(chat_url, json=payload, headers=headers, timeout=3600)
        response.raise_for_status()  # 如果请求失败则抛出异常
        
        response_data = response.json()
        content = response_data['choices'][0]['message']['content'].lstrip()
        
        # 返回与其它query函数格式一致的结果
        features = {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}
        return features, content

    except requests.exceptions.RequestException as e:
        error_msg = f"vLLM request failed: {str(e)}"
        console.print(f"[bold red]{error_msg}")
        features = {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}
        return features, f"[bold red]{error_msg}"

def query_llm_online(api_url, api_key, payload: Dict[str, str], tokenizer_path: str)-> str:
    """Call SiliconFlow API to get LLM response"""
    import requests
    from rich.console import Console
    # 将messages转换成字符串格式
    conversation= ""
    for message in payload["messages"]:
        conversation+= f"{message['role']}: {message['content']}"
        
    console= Console()
    headers= {
        "Authorization": api_key,
        "Content-Type": "application/json"
    }
    try:
        response= requests.post(
            api_url,
            json= payload,
            headers= headers,
            timeout= 3600
        )
        response.raise_for_status()
        return {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}, response.json()['choices'][0]['message']['content'].lstrip()
    except Exception as e:
        console.print(f"[bold red]API request failed: {str(e)}")
        return {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}, f"[bold red]API request failed: {str(e)}"

def query_llm(model:str, model_folder: str, messages:List, temperature:float= 0.6, max_token:int= 1024, top_p:float= 0.9, repetition_penalty:float= 1.1):
    from transformers import AutoTokenizer, AutoModelForCausalLM
    import torch
    # 加载预训练模型和分词器
    tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
    tokenizer = AutoTokenizer.from_pretrained(tokenizer_path, device_map= "cuda")
    local_model = AutoModelForCausalLM.from_pretrained(
        os.path.join(model_folder, model),
        torch_dtype="float16",
        low_cpu_mem_usage=True,
        device_map="cuda", offload_state_dict= False,
        )

    # 将messages转换成字符串格式
    conversation = ""
    for message in messages:
        conversation += f"{message['role']}: {message['content']}"

    # 编码输入
    input_ids = tokenizer.encode(conversation + tokenizer.eos_token, return_tensors='pt').to("cuda")

    # 生成回复
    output = local_model.generate(
        input_ids, 
        pad_token_id=tokenizer.eos_token_id,
        max_new_tokens= max_token,
        temperature= temperature, 
        top_p= top_p,
        repetition_penalty= repetition_penalty
    )
    response = tokenizer.decode(output[:, input_ids.shape[-1]:][0], skip_special_tokens=True)

    del local_model
    del tokenizer
    gc.collect()
    torch.cuda.empty_cache()
    return {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}, response

@cpu(cpu_num= 1, mem= 1024)
def task1_file_process(context):
    """
    Processes the file content based on type and stores the processed content in context.
    """
    try:
        start_time= time.time()
        dag_id= ray.get(context.get.remote('dag_id'))
        question = ray.get(context.get.remote("question"))
        supplementary_files = ray.get(context.get.remote("supplementary_files"))
        file_name = ""
        content = None
        
        if supplementary_files:
            file_name = next(iter(supplementary_files.keys()), "")
            content = next(iter(supplementary_files.values()), None)
            print(f"  -> Found single supplementary file to process: '{file_name}'")
        else:
            print("  -> No supplementary files found in context.")
        processed_content = ""
        
        print("✅ Task 1: Starting file processing...")
        if content is None:
            print(f"  -> No content for file '{file_name}', skipping processing.")
        else:
            print(f"  -> Processing file: {file_name}")
            file_name_lower = file_name.lower()
            
            # Excel 文件处理
            if file_name_lower.endswith(('.xlsx', '.xls')):
                print("处理 Excel 文件...")
                excel_file = BytesIO(content)
                
                # 读取所有sheet
                excel_data = []
                with pd.ExcelFile(excel_file) as xls:
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(xls, sheet_name)
                        excel_data.append(f"=== Sheet: {sheet_name} ===\n")
                        excel_data.append(df.to_string(index=True))
                
                processed_content = "\n\n".join(excel_data)
                print(f"提取到 {len(excel_data)//2} 个工作表的内容")

            # CSV 文件处理
            elif file_name_lower.endswith('.csv'):
                print("处理 CSV 文件...")
                csv_file = BytesIO(content)
                df = pd.read_csv(csv_file)
                processed_content = df.to_string(index=True)
                print(f"提取到 CSV 内容，共 {len(df)} 行")

            # XML 文件处理
            elif file_name_lower.endswith('.xml'):
                print("[File Process] 处理 XML 文件...")
                try:
                    xml_content = content.decode('utf-8')
                    print("[File Process] XML 内容解码成功")
                    
                    # 尝试解析和格式化 XML
                    try:
                        parsed_xml = minidom.parseString(xml_content)
                        processed_content = parsed_xml.toprettyxml()
                        print("[File Process] XML 格式化成功")
                    except Exception as xml_error:
                        print(f"[File Process] XML 格式化失败: {str(xml_error)}")
                        processed_content = xml_content
                    
                    # 提取文本内容（去除 XML 标签）
                    try:
                        root = ET.fromstring(xml_content)
                        def extract_text(element):
                            text = element.text or ''
                            for child in element:
                                text += extract_text(child)
                                if child.tail:
                                    text += child.tail
                            return text
                        
                        text_content = extract_text(root)
                        print("[File Process] 文本内容提取成功")
                        
                        processed_content = (
                            "=== XML Structure ===\n"
                            f"{processed_content}\n\n"
                            "=== Extracted Text Content ===\n"
                            f"{text_content}"
                        )
                    except Exception as text_error:
                        print(f"[File Process] 文本提取失败: {str(text_error)}")
                    
                    print(f"[File Process] XML 处理完成，内容长度: {len(processed_content)}")
                    
                except Exception as e:
                    print(f"[File Process] XML 处理出错: {str(e)}")
                    processed_content = content.decode('utf-8')

            # PPT 文件处理
            elif file_name_lower.endswith('.pptx'):
                print("处理 PPT 文件...")
                prs = Presentation(BytesIO(content))
                all_slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text_parts.append(shape.text.strip())
                    all_slides_text.append(f"--- Slide {i+1} ---\n" + '\n'.join(filter(None, slide_text_parts)))
                processed_content = "\n\n".join(all_slides_text)
                print(f"提取到 {len(all_slides_text)} 页幻灯片内容")

            # Word 文件处理
            elif file_name_lower.endswith('.docx'):
                print("处理 Word 文件...")
                doc = Document(BytesIO(content))
                paragraphs_text = []
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs_text.append(text)
                processed_content = "\n".join(paragraphs_text)
                print(f"提取到 {len(paragraphs_text)} 段文本")

            # PDF 文件处理
            elif file_name_lower.endswith('.pdf'):
                print("处理 PDF 文件...")
                temp_pdf_path = f"/tmp/temp_{file_name_lower.split('.')[0]}.pdf"
                with open(temp_pdf_path, "wb") as f:
                    f.write(content)
                
                pdf_pages = []
                with pdfplumber.open(temp_pdf_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            pdf_pages.append(f"--- Page {i+1} ---\n{page_text}")
                        
                        tables = page.extract_tables()
                        if tables:
                            table_texts = []
                            for j, table in enumerate(tables):
                                table_text = f"Table {j+1} on Page {i+1}:\n"
                                for row in table:
                                    cleaned_row = [str(cell) if cell is not None else "" for cell in row]
                                    table_text += " | ".join(cleaned_row) + "\n"
                                table_texts.append(table_text)
                            
                            if table_texts:
                                pdf_pages.append("\n".join(table_texts))
                
                os.remove(temp_pdf_path)
                processed_content = "\n\n".join(pdf_pages)
                print(f"提取到 {len(pdf_pages)} 页/节 PDF 内容")

            # 文本文件处理
            elif file_name_lower.endswith('.txt'):
                print("处理 TXT 文件...")
                processed_content = content.decode('utf-8')

            else:
                print(f"警告: 未知或不支持的文件类型: {file_name_lower}，尝试按文本处理")
                try:
                    processed_content = content.decode('utf-8')
                except UnicodeDecodeError:
                    print("文件无法解码为 UTF-8，将跳过内容注入")
                    processed_content = "[Content could not be decoded]"
            print("✅ Task 1: File processing finished.")
        # 生成推理模型的特征
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        llm_process_feat= {"text_length": len(prompt), "token_count": estimate_tokens(prompt)}
        # 存储处理后的内容
        context.put.remote("processed_content", processed_content)
        print(f"✅ 文件处理结束...")
        return json.dumps({ # 为了保存到本地
            "dag_id": dag_id,
            "succ_task_feat": {
                "task2_llm_process_qwen": {"text_length": llm_process_feat["text_length"], "token_count": llm_process_feat["token_count"], "reason": 1},
                "task3_llm_process_deepseek": {"text_length": llm_process_feat["text_length"], "token_count": llm_process_feat["token_count"], "reason": 0}
            },
            "curr_task_feat": None,
            "start_time": start_time,
            "end_time": time.time()
        })
    except Exception as e:
        print(f"task2_file_process 发生错误: {str(e)}")
        raise e

@gpu(gpu_mem= 70000, model_name= "qwen3-32b", backend="huggingface")
def task2_llm_process_qwen(context):
    
    """
    Processes the file content using LLM based on the question.
    """
    try:
        backend= task2_llm_process_qwen._task_decorator['backend']
        print(f"✅ qwen处理开始....")
        start_time= time.time()  # ADD: 记录开始时间 
        dag_id= ray.get(context.get.remote("dag_id"))
        processed_content = ray.get(context.get.remote("processed_content")) 
        question= ray.get(context.get.remote("question"))

        use_online_model= ray.get(context.get.remote("use_online_model"))
        model_folder= ray.get(context.get.remote("model_folder"))
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        temperature, max_token, repetition_penalty, top_p= None, None, None, None
        if not use_online_model:
            temperature= ray.get(context.get.remote("temperature"))
            max_token= ray.get(context.get.remote("max_tokens"))
            top_p= ray.get(context.get.remote("top_p"))
            repetition_penalty= ray.get(context.get.remote("repetition_penalty"))
        if not question:
             raise ValueError(f"任务 {dag_id} 缺少 Question")

        # 构建提示
        # 生成推理模型的特征
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        payload= {
            "model": "Qwen/Qwen3-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": temperature,
            "max_tokens": max_token,
            "enable_thinking": False,
            "response_format": {"type": "text"}               
        }
        inference_features, answer= None, None
        if use_online_model:
            inference_features, answer= query_llm_online(
                api_url= ray.get(context.get.remote("api_url")), 
                api_key= ray.get(context.get.remote("api_key")), 
                payload= payload,
                tokenizer_path= tokenizer_path)
        elif backend == "vllm":
            inference_features, answer= query_vllm_model(
                api_url= ray.get(context.get.remote(f"task2_llm_process_qwen_request_api_url")), 
                model_alias= "qwen3-32b",
                messages= payload["messages"],
                temperature= temperature,
                max_token= max_token,
                top_p= top_p,
                repetition_penalty= repetition_penalty)
        else:
            inference_features, answer= query_llm(model_folder= model_folder, model= "Qwen/Qwen3-32B", messages= [{"role": "user", "content": prompt}], temperature= temperature, max_token= max_token, top_p= top_p, repetition_penalty= repetition_penalty)

        inference_features['reason'] = 0
        context.put.remote("qwen_answer", answer)
        next_task_prompt = (
            "You are a senior editor and a world-class reasoning expert. Your job is to synthesize the answers from two different AI assistants to produce one final, superior answer for the given question.\n\n"
            f"--- Original Question ---\n{question}\n\n"
            f"--- Answer from Assistant 1 (Qwen3) ---"
            f"--- Answer from Assistant 2 (DeepSeek) ---"
            "--- Your Task ---\n"
            "Analyze both answers. Identify the strengths and weaknesses of each. Then, combine their best elements, correct any errors, and provide a single, comprehensive, and accurate final answer. Adhere to the final answer format requested in the original prompt.\n\n"
            "Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]."
        )
        text1_length= len(answer)
        text1_token_count= estimate_tokens(answer)
        context.put.remote("text1_feature", {"text1_length": text1_length, "text1_token_count": text1_token_count})
        return json.dumps({
            "dag_id": dag_id,
            "curr_task_feat": inference_features,
            "succ_task_feat": {
                "task4_llm_fuse_answer":{
                    "prompt_length": len(next_task_prompt),
                    "prompt_token_count": estimate_tokens(next_task_prompt),
                    "text1_length": text1_length,
                    "text1_token_count": text1_token_count,
                    "reason": 0
                }
            },
            "start_time": start_time,
            "end_time": time.time()
        })
    except Exception as e:
        print(f"task2_llm_process_qwen 发生错误: {str(e)}")
        raise e

@gpu(gpu_mem= 70000, model_name= "deepseek-r1-32b", backend="huggingface")
def task3_llm_process_deepseek(context):
    
    """
    Processes the file content using LLM based on the question.
    """
    try:
        backend= task3_llm_process_deepseek._task_decorator['backend']
        start_time= time.time()  # ADD: 记录开始时间 
        print(f"✅ deepseek处理开始....")
        dag_id= ray.get(context.get.remote("dag_id"))
        processed_content = ray.get(context.get.remote("processed_content")) 
        question= ray.get(context.get.remote("question"))

        use_online_model= ray.get(context.get.remote("use_online_model"))
        model_folder= ray.get(context.get.remote("model_folder"))
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        temperature, max_token, repetition_penalty, top_p= None, None, None, None
        if not use_online_model:
            temperature= ray.get(context.get.remote("temperature"))
            max_token= ray.get(context.get.remote("max_tokens"))
            top_p= ray.get(context.get.remote("top_p"))
            repetition_penalty= ray.get(context.get.remote("repetition_penalty"))
        if not question:
             raise ValueError(f"任务 {dag_id} 缺少 Question")
        # 构建提示
        # 生成推理模型的特征
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your concise thinking thoughts and don't think too complicated, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        payload= {
            "model": "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": temperature,
            "max_tokens": max_token,
            "response_format": {"type": "text"}
        }
        inference_features, answer= None, None
        if use_online_model:
            inference_features, answer= query_llm_online(
                api_url= ray.get(context.get.remote("api_url")), 
                api_key= ray.get(context.get.remote("api_key")), 
                payload= payload,
                tokenizer_path= tokenizer_path)
        elif backend == "vllm":
            inference_features, answer= query_vllm_model(
                api_url= ray.get(context.get.remote(f"task3_llm_process_deepseek_request_api_url")),
                model_alias= "deepseek-r1-32b",
                messages= payload["messages"],
                temperature= temperature,
                max_token= max_token,
                top_p= top_p,
                repetition_penalty= repetition_penalty)
        else:
            inference_features, answer= query_llm(model_folder= model_folder, model= "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B", messages= [{"role": "user", "content": prompt}], temperature= temperature, max_token= max_token, top_p= top_p, repetition_penalty= repetition_penalty)        
        
        context.put.remote("deepseek_answer", answer)
        inference_features['reason'] = 1
        next_task_prompt = (
            "You are a senior editor and a world-class reasoning expert. Your job is to synthesize the answers from two different AI assistants to produce one final, superior answer for the given question.\n\n"
            f"--- Original Question ---\n{question}\n\n"
            f"--- Answer from Assistant 1 (Qwen3) ---"
            f"--- Answer from Assistant 2 (DeepSeek) ---"
            "--- Your Task ---\n"
            "Analyze both answers. Identify the strengths and weaknesses of each. Then, combine their best elements, correct any errors, and provide a single, comprehensive, and accurate final answer. Adhere to the final answer format requested in the original prompt.\n\n"
            "Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]."
        )
        text2_length= len(answer)
        text2_token_count= estimate_tokens(answer)
        context.put.remote("text2_feature", {"text2_length": text2_length, "text2_token_count": text2_token_count})
        return json.dumps({
            "dag_id": dag_id,
            "curr_task_feat": inference_features,
            "succ_task_feat": {
                "task4_llm_fuse_answer":{
                    "prompt_length": len(next_task_prompt),
                    "prompt_token_count": estimate_tokens(next_task_prompt),         
                    "text2_length": text2_length,
                    "text2_token_count": text2_token_count,
                    "reason": 0
                }
            },            
            "start_time": start_time,
            "end_time": time.time()
        })
    except Exception as e:
        print(f"task3_llm_process_deepseek 发生错误: {str(e)}")
        raise e

@gpu(gpu_mem= 70000, model_name= "qwen3-32b", backend="huggingface") # 融合任务通常也需要LLM，但可能比生成任务资源消耗小
def task4_llm_fuse_answer(context):
    """
    Fusion Task: Takes answers from multiple experts, analyzes them, and generates a final, synthesized answer.
    """
    try:
        backend= task4_llm_fuse_answer._task_decorator['backend']
        print("✅ Task 4 (Fusion): Starting answer synthesis...")
        start_time= time.time()  # ADD: 记录开始时间 
        # 1. 从 context 获取所有专家的答案
        answer_qwen = ray.get(context.get.remote("qwen_answer"))
        answer_deepseek = ray.get(context.get.remote("deepseek_answer"))
        dag_id = ray.get(context.get.remote("dag_id"))
        question = ray.get(context.get.remote("question"))
        text1_feature= ray.get(context.get.remote("text1_feature"))
        text2_feature= ray.get(context.get.remote("text2_feature"))
        use_online_model= ray.get(context.get.remote("use_online_model"))
        model_folder= ray.get(context.get.remote("model_folder"))
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        temperature, max_token, repetition_penalty, top_p= None, None, None, None
        if not use_online_model:
            temperature= ray.get(context.get.remote("temperature"))
            max_token= ray.get(context.get.remote("max_tokens"))
            top_p= ray.get(context.get.remote("top_p"))
            repetition_penalty= ray.get(context.get.remote("repetition_penalty"))
        if not question:
             raise ValueError(f"任务 {dag_id} 缺少 Question")

        # 2. 构建一个用于融合的 prompt
        prompt = (
            "You are a senior editor and a world-class reasoning expert. Your job is to synthesize the answers from two different AI assistants to produce one final, superior answer for the given question.\n\n"
            f"--- Original Question ---\n{question}\n\n"
            f"--- Answer from Assistant 1 (Qwen3) ---\n{answer_qwen}\n\n"
            f"--- Answer from Assistant 2 (DeepSeek) ---\n{answer_deepseek}\n\n"
            "--- Your Task ---\n"
            "Analyze both answers. Identify the strengths and weaknesses of each. Then, combine their best elements, correct any errors, and provide a single, comprehensive, and accurate final answer. Adhere to the final answer format requested in the original prompt.\n\n"
            "Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]."
        )
        
        # 3. 使用一个强大的模型（如Qwen3）来执行融合
        payload= {
            "model": "Qwen/Qwen3-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": temperature,
            "max_tokens": max_token,
            "enable_thinking": False,
            "response_format": {"type": "text"}               
        }
        final_answer= None
        if use_online_model:
            inference_features, final_answer= query_llm_online(
                api_url= ray.get(context.get.remote("api_url")), 
                api_key= ray.get(context.get.remote("api_key")), 
                payload= payload,
                tokenizer_path= tokenizer_path)
        elif backend == "vllm":
            inference_features, final_answer= query_vllm_model(
                api_url= ray.get(context.get.remote(f"task4_llm_fuse_answer_request_api_url")),
                model_alias= "qwen3-32b",
                messages= payload["messages"],
                temperature= temperature,
                max_token= max_token,
                top_p= top_p,
                repetition_penalty= repetition_penalty)
        else:
            inference_features, final_answer= query_llm(model_folder= model_folder, model= "Qwen/Qwen3-32B", messages= [{"role": "user", "content": prompt}], temperature= temperature, max_token= max_token, top_p= top_p, repetition_penalty= repetition_penalty)        
        # 4. 准备并返回最终的输出
        print(f"  -> Fusion complete. Final answer generated for DAG {dag_id}.")
        return json.dumps({
            "dag_id": dag_id,
            "final_answer": final_answer,
            "curr_task_feat": 
            {
                "prompt_length": inference_features["text_length"], 
                "prompt_token_count": inference_features["token_count"], 
                "text1_length": text1_feature["text1_length"],
                "text1_token_count": text1_feature["text1_token_count"],
                "text2_length": text2_feature["text2_length"],
                "text2_token_count": text2_feature["text2_token_count"],
                "reason": 0
            },
            "start_time": start_time,
            "end_time": time.time()
        })
    except Exception as e:
        print(f"❌ task4_fuse_final_answer failed: {str(e)}")
        raise e